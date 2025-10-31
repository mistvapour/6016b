#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
预答辩PPT生成脚本
基于MIL-STD-6016的战术数据链信息标准数据库架构设计与整合应用
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("需要安装python-pptx库")
    print("请运行: pip install python-pptx")
    exit(1)

def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 第1页：标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题页布局
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "基于MIL-STD-6016的战术数据链信息标准数据库架构设计与整合应用"
    subtitle.text = """Database Architecture Design and Integration Application of 
Tactical Data Link Information Standards Based on MIL-STD-6016

答辩人：[您的姓名]
指导教师：[导师姓名]
专业：计算机科学与技术
答辩时间：2025年X月X日
华东师范大学"""
    
    # 第2页：研究背景与意义
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "研究背景与意义"
    content2.text = """研究背景
• 信息化战争发展：战术数据链（TDL）成为现代作战体系核心
• Link16广泛应用：基于MIL-STD-6016标准，实现跨平台信息交互
• 多链路挑战：信号检测困难、天基拓展需求、多链融合挑战

研究意义
• 理论价值：推动不同数据链协议的统一
• 工程意义：为装备间互操作提供技术基础
• 应用价值：支持多链融合与仿真验证"""
    
    # 第3页：国内外研究现状
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "国内外研究现状"
    content3.text = """国外研究
• 标准化体系：MIL-STD-6016、MIL-STD-3011、STANAG 5516
• 技术支撑：MITRE多链融合与互操作工程
• 性能优化：抗干扰建模、态势信息处理

国内研究
• 信号检测：Link16信号检测与识别方法
• 天基拓展：天基拓展及其影响研究
• 仿真平台：数据链仿真平台构建

研究空白
• 标准数据库设计与多链融合整合方面有待深入"""
    
    # 第4页：研究内容与目标
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "研究内容与目标"
    content4.text = """主要研究内容
1. 标准梳理：系统梳理Link16消息标准，构建标准化数据库模型
2. 架构设计：设计数据库架构，支持消息存储、查询与一致性校验
3. 系统对接：实现数据库与仿真系统的对接
4. 性能评估：通过实验仿真评估数据库效果

研究目标
• 实现J系列报文结构、字段及编码规则的统一管理
• 支持跨消息字段查询以及数据一致性校验
• 验证数据库在态势信息处理和多链融合应用中的可行性"""
    
    # 第5页：研究方法与技术路线
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "研究方法与技术路线"
    content5.text = """研究方法
1. 文献调研与标准分析：研读相关研究成果与标准文档
2. 数据库建模与架构设计：构建信息标准数据库模型
3. 系统实现与整合应用：开发数据库应用接口
4. 实验仿真与性能验证：利用仿真平台进行实验评估

技术路线
标准分析 → 需求分析 → 数据库设计 → 系统实现 → 测试验证"""
    
    # 第6页：系统需求分析
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "系统需求分析"
    content6.text = """功能需求
• 标准消息管理：J系列报文的存储、查询与管理
• 字段与语义概念绑定：字段与语义概念的关联
• 多链路互操作支持：跨链路消息对接与映射

性能需求
• 查询效率：支持大规模数据的高效检索
• 并发处理：支持多用户同时访问
• 数据一致性：保证数据的完整性和一致性

安全需求
• 权限管理：基于角色的访问控制
• 数据保护：敏感信息的安全存储"""
    
    # 第7页：数据库总体设计
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "数据库总体设计"
    content7.text = """设计目标
• 标准化存储：遵循MIL-STD-6016和STANAG 5516标准
• 语义绑定：实现字段与语义概念的绑定
• 互操作支持：支持跨链路数据互操作
• 高性能检索：提供高效的索引与查询机制

设计原则
• 模块化：分为消息、字段、概念、映射等子模块
• 分层化：逻辑层、物理层与接口层分离
• 面向应用：充分考虑实际使用场景"""
    
    # 第8页：数据库架构设计
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "数据库架构设计"
    content8.text = """四层架构
1. 数据存储层：关系数据库模式，存储核心实体
2. 数据管理层：提供数据解析、校验、绑定功能
3. 接口服务层：通过RESTful API与外部系统交互
4. 应用展示层：提供前端可视化和交互操作

核心实体
• MESSAGE：J报文元数据，带标准版本维度
• FIELD：位段定义，约束位段不重叠
• CONCEPT：语义概念库
• MAPPING：跨链/跨版映射规则与置信度"""
    
    # 第9页：数据表结构设计
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局，用于表格
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title9.text_frame.text = "数据表结构设计"
    title9.text_frame.paragraphs[0].font.size = Pt(24)
    title9.text_frame.paragraphs[0].font.bold = True
    
    # 添加表格
    table = slide9.shapes.add_table(5, 4, Inches(0.5), Inches(2), Inches(12), Inches(4)).table
    
    # 设置表格标题行
    table.cell(0, 0).text = "表名"
    table.cell(0, 1).text = "主键/唯一约束"
    table.cell(0, 2).text = "关键字段"
    table.cell(0, 3).text = "说明"
    
    # 设置表格内容
    table.cell(1, 0).text = "MESSAGE"
    table.cell(1, 1).text = "PK(message_id)"
    table.cell(1, 2).text = "j_num, j_series, title"
    table.cell(1, 3).text = "J报文元数据"
    
    table.cell(2, 0).text = "FIELD"
    table.cell(2, 1).text = "PK(field_id)"
    table.cell(2, 2).text = "start_bit, end_bit, unit"
    table.cell(2, 3).text = "位段定义"
    
    table.cell(3, 0).text = "CONCEPT"
    table.cell(3, 1).text = "PK(concept_id)"
    table.cell(3, 2).text = "name, definition, source"
    table.cell(3, 3).text = "语义概念库"
    
    table.cell(4, 0).text = "MAPPING"
    table.cell(4, 1).text = "PK(map_id)"
    table.cell(4, 2).text = "message_id, field_id, rule"
    table.cell(4, 3).text = "跨链映射规则"
    
    # 第10页：系统架构设计
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "系统架构设计"
    content10.text = """总体架构
• 数据层：MySQL数据库，存储核心数据
• 服务层：FastAPI实现RESTful接口
• 应用层：React前端交互模块
• 外部接口层：与仿真平台、指挥系统对接

技术栈
• 后端：FastAPI + MySQL 8.0
• 前端：React + shadcn/ui
• 部署：Docker容器化部署"""
    
    # 第11页：核心功能实现
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "核心功能实现"
    content11.text = """数据导入功能
• CSV/Excel批量导入：支持大规模数据自动导入
• 数据清洗：自动处理异常数据和格式转换
• 数据校验：确保导入数据的完整性和一致性

查询检索功能
• 多条件搜索：支持关键词、J系列、模糊/精确搜索
• 跨标准比较：支持不同版本规范的对比分析
• 语义绑定：实现字段与概念的双向绑定"""
    
    # 第12页：系统实现展示
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "系统实现展示"
    content12.text = """主要功能模块
1. 消息管理：J系列报文的增删改查
2. 字段管理：位段信息的结构化存储
3. 概念管理：语义概念的定义和维护
4. 映射管理：跨标准映射关系的建立

用户界面
• 搜索界面：提供直观的搜索和筛选功能
• 结果展示：表格化展示搜索结果
• 统计分析：提供数据统计和可视化图表"""
    
    # 第13页：系统测试方案
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    title13 = slide13.shapes.title
    content13 = slide13.placeholders[1]
    
    title13.text = "系统测试方案"
    content13.text = """测试目标
• 数据库正确性：验证数据完整性和一致性
• 接口稳定性：测试API的并发性能
• 功能可用性：验证用户界面的易用性
• 安全鲁棒性：测试系统的安全防护能力

测试方法
• 功能测试：手工操作与自动化脚本
• 压力测试：Apache JMeter模拟高并发
• 安全测试：SQL注入、参数验证等"""
    
    # 第14页：性能测试结果
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    title14 = slide14.shapes.title
    content14 = slide14.placeholders[1]
    
    title14.text = "性能测试结果"
    content14.text = """数据库性能
• 数据规模：支持50万条以上消息字段
• 查询效率：平均响应延迟小于150ms
• 并发处理：支持1000并发用户访问

接口性能
• 搜索接口：100并发下延迟<150ms
• 比较接口：跨5个规范版本响应约200ms
• 成功率：1000并发下保持95%成功率"""
    
    # 第15页：功能测试结果
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    title15 = slide15.shapes.title
    content15 = slide15.placeholders[1]
    
    title15.text = "功能测试结果"
    content15.text = """数据库层面
• 约束验证：外键约束、唯一约束正常工作
• 数据一致性：位宽一致性检查通过
• 审计功能：异常记录自动记录至审计表

接口层面
• 参数验证：非法参数正确返回错误
• 安全防护：SQL注入等攻击有效防护
• 权限控制：RBAC角色分离策略有效

前端层面
• 兼容性：多浏览器兼容性良好
• 用户体验：搜索条件正确回显
• 错误处理：弱网条件下降级机制正常"""
    
    # 第16页：系统特色与创新点
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    title16 = slide16.shapes.title
    content16 = slide16.placeholders[1]
    
    title16.text = "系统特色与创新点"
    content16.text = """主要特色
1. 标准化建模：基于MIL-STD-6016的规范化数据模型
2. 语义绑定：字段与概念的双向绑定机制
3. 跨链支持：支持多链路互操作和映射
4. 高性能：优化的索引和查询策略

创新点
• 统一数据模型：首次系统化建模MIL-STD-6016标准
• 语义一致性：通过概念绑定提升跨标准一致性
• 工程化实现：完整的系统架构和接口设计"""
    
    # 第17页：应用价值与前景
    slide17 = prs.slides.add_slide(prs.slide_layouts[1])
    title17 = slide17.shapes.title
    content17 = slide17.placeholders[1]
    
    title17.text = "应用价值与前景"
    content17.text = """直接应用
• 标准管理：为MIL-STD-6016标准提供数字化管理平台
• 研发支持：为数据链系统开发提供标准参考
• 培训教育：为相关专业教学提供实践平台

扩展应用
• 多链融合：支持Link16、JREAP、TTNT等多链路整合
• 仿真验证：为战术数据链仿真提供数据支撑
• 互操作测试：为装备互操作性验证提供工具"""
    
    # 第18页：存在的问题与不足
    slide18 = prs.slides.add_slide(prs.slide_layouts[1])
    title18 = slide18.shapes.title
    content18 = slide18.placeholders[1]
    
    title18.text = "存在的问题与不足"
    content18.text = """数据规模限制
• 样本有限：实验数据约50万条，未覆盖全部NATO规范
• 跨链验证：缺乏Link11/22、TTNT的大规模验证

功能深度
• 可视化：前端以表格展示为主，缺乏复杂图表分析
• 冲突检测：一致性验证逻辑仍偏基础

安全运维
• 加密保护：未覆盖全链路加密、KMI等高级安全要求
• 运维监控：缺乏完整的运维监控体系"""
    
    # 第19页：未来研究方向
    slide19 = prs.slides.add_slide(prs.slide_layouts[1])
    title19 = slide19.shapes.title
    content19 = slide19.placeholders[1]
    
    title19.text = "未来研究方向"
    content19.text = """短期目标
1. 数据扩展：导入更多NATO标准数据
2. 功能增强：增加复杂查询和可视化功能
3. 性能优化：进一步提升查询效率和并发能力

长期目标
1. 多链融合：支持更多战术数据链标准
2. 智能分析：引入AI技术进行智能分析
3. 分布式部署：支持大规模分布式部署"""
    
    # 第20页：总结与致谢
    slide20 = prs.slides.add_slide(prs.slide_layouts[1])
    title20 = slide20.shapes.title
    content20 = slide20.placeholders[1]
    
    title20.text = "总结与致谢"
    content20.text = """主要贡献
• 理论贡献：建立了MIL-STD-6016的标准化数据模型
• 技术贡献：实现了高性能的数据库架构和接口设计
• 应用贡献：为战术数据链互操作提供了技术支撑

致谢
感谢导师的悉心指导！
感谢实验室同学的支持！
感谢各位专家的宝贵意见！

请各位专家批评指正！"""
    
    return prs

def main():
    """主函数"""
    try:
        print("正在创建预答辩PPT...")
        prs = create_presentation()
        
        # 保存PPT文件
        filename = "基于MIL-STD-6016的战术数据链信息标准数据库架构设计与整合应用_预答辩.pptx"
        prs.save(filename)
        
        print(f"PPT创建成功！文件保存为：{filename}")
        print("PPT包含20页内容，涵盖了论文的所有核心要点")
        
    except Exception as e:
        print(f"创建PPT时出错：{e}")
        print("请确保已安装python-pptx库：pip install python-pptx")

if __name__ == "__main__":
    main()

